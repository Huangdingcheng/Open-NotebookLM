from __future__ import annotations
import os
import re
import asyncio
from pathlib import Path
from typing import List, Dict, Any, Optional

import fitz  # PyMuPDF
from workflow_engine.workflow.registry import register
from workflow_engine.graphbuilder.graph_builder import GenericGraphBuilder
from workflow_engine.logger import get_logger
from workflow_engine.state import IntelligentQAState, MainState
from workflow_engine.agentroles import create_vlm_agent, create_agent
from workflow_engine.utils import get_project_root
from workflow_engine.promptstemplates.resources.pt_qa_agent_repo import QaAgent as QaAgentPrompts

log = get_logger(__name__)

# 文档上下文长度限制（字符数）
MAX_DOC_CONTEXT_CHARS = 48000
RAG_TOP_K = 30
MAX_HISTORY_TURNS = 10

# Try importing office libraries
try:
    from docx import Document
except ImportError:
    Document = None

try:
    from pptx import Presentation
except ImportError:
    Presentation = None

@register("intelligent_qa")
def create_intelligent_qa_graph() -> GenericGraphBuilder:
    """
    Workflow for Intelligent Q&A with Parallel Processing
    Steps:
    1. Parse uploaded files (PDF/Office/Image/Video)
    2. Parallel Analysis: Call LLM/VLM for each file individually
    3. Aggregate context and Final QA
    """
    builder = GenericGraphBuilder(state_model=IntelligentQAState, entry_point="_start_")

    def _extract_text_result(state: MainState, role_name: str) -> str:
        try:
            result = state.agent_results.get(role_name, {}).get("results", {})
            if isinstance(result, dict):
                return result.get("text") or result.get("raw") or ""
            if isinstance(result, str):
                return result
        except Exception:
            return ""
        return ""

    def _get_embedded_file_paths(state: IntelligentQAState) -> set:
        """从 vector store manifest 读取已入库文件的 resolved 路径集合。"""
        base_dir = getattr(state.request, "vector_store_base_dir", None) or ""
        if not base_dir:
            return set()
        try:
            from workflow_engine.toolkits.ragtool.vector_store_tool import VectorStoreManager
            manager = VectorStoreManager(base_dir=base_dir)
            manifest_files = manager.manifest.get("files", []) or []
            embedded = set()
            for f in manifest_files:
                orig = f.get("original_path") or ""
                if orig:
                    try:
                        embedded.add(str(Path(orig).resolve()))
                    except Exception:
                        pass
            return embedded
        except Exception as e:
            log.debug(f"Could not read embedded file paths: {e}")
            return set()

    def _start_(state: IntelligentQAState) -> IntelligentQAState:
        # Ensure request fields
        if not state.request.file_ids:
            state.request.file_ids = []
        if not state.request.query:
            state.request.query = ""
        # Initialize file analyses
        state.file_analyses = []
        return state

    async def parallel_parse_node(state: IntelligentQAState) -> IntelligentQAState:
        """
        Parallel parsing AND analysis of all files
        """
        files = state.request.file_ids
        if not files:
            state.context_content = ""
            return state
        
        def _infer_target_files(query: str, file_paths: List[str]) -> List[str]:
            if not query:
                return []
            q = query.lower()
            q_compact = re.sub(r"\s+", "", q)
            matches: List[str] = []
            for path in file_paths:
                name = Path(path).name.lower()
                stem = Path(path).stem.lower()
                name_compact = re.sub(r"\s+", "", name)
                stem_compact = re.sub(r"\s+", "", stem)
                if name in q or stem in q or name_compact in q_compact or stem_compact in q_compact:
                    matches.append(path)
            return matches

        target_files = _infer_target_files(state.request.query or "", files)
        files_to_process = target_files if target_files else files

        # Filter out already-embedded files — they will be handled by RAG retrieval in chat_node
        embedded_paths = _get_embedded_file_paths(state)
        if embedded_paths:
            before_count = len(files_to_process)
            files_to_process = [
                f for f in files_to_process
                if str(Path(f).resolve()) not in embedded_paths
            ]
            skipped = before_count - len(files_to_process)
            if skipped:
                log.info(f"Skipping {skipped} embedded files from parallel parse")

        async def process_file(file_path: str) -> Dict[str, Any]:
            file_path_obj = Path(file_path)
            filename = file_path_obj.name
            
            if not file_path_obj.exists():
                return {
                    "filename": filename,
                    "analysis": f"[Error: File not found {file_path}]",
                    "content": ""
                }
            
            suffix = file_path_obj.suffix.lower()
            raw_content = ""
            analysis_result = ""
            file_type = "unknown"
            
            try:
                # ==========================
                # 1. Extraction Phase
                # ==========================
                
                # PDF
                if suffix == ".pdf":
                    file_type = "document"
                    try:
                        doc = fitz.open(file_path)
                        text = ""
                        for page in doc:
                            text += page.get_text() + "\n"
                        raw_content = text
                    except Exception as e:
                        raw_content = f"[Error parsing PDF: {e}]"

                # Word
                elif suffix in [".docx", ".doc"]:
                    file_type = "document"
                    if Document is None:
                         raw_content = "[Error: python-docx not installed]"
                    else:
                        try:
                            doc = Document(file_path)
                            raw_content = "\n".join([p.text for p in doc.paragraphs])
                        except Exception as e:
                             raw_content = f"[Error parsing Docx: {e}]"

                # PPT (Simulate PPT -> Text)
                elif suffix in [".pptx", ".ppt"]:
                    file_type = "presentation"
                    if Presentation is None:
                        raw_content = "[Error: python-pptx not installed]"
                    else:
                        try:
                            prs = Presentation(file_path)
                            text = ""
                            for i, slide in enumerate(prs.slides):
                                text += f"--- Slide {i+1} ---\n"
                                for shape in slide.shapes:
                                    if hasattr(shape, "text"):
                                        text += shape.text + "\n"
                            raw_content = text
                        except Exception as e:
                            raw_content = f"[Error parsing PPT: {e}]"
                
                # Image / Video
                elif suffix in [".jpg", ".jpeg", ".png", ".mp4", ".mov", ".avi"]:
                    file_type = "media"
                    # For media, raw_content comes from VLM understanding
                    # We will do analysis directly here using VLM Agent
                    raw_content = "[Media file - will be analyzed by VLM]"
                
                else:
                    file_type = "text"
                    try:
                        with open(file_path, "r", encoding="utf-8") as f:
                            raw_content = f.read()
                    except:
                        raw_content = "[Unsupported file type]"

                # ==========================
                # 2. Analysis Phase (Parallel LLM Call)
                # ==========================

                if file_type == "media":
                    # Use VLM Agent for Media
                    vlm_mode = "understanding"
                    input_key = "input_image"
                    if suffix in [".mp4", ".mov", ".avi"]:
                        vlm_mode = "video_understanding"
                        input_key = "input_video"

                    try:
                        vlm_prompt = QaAgentPrompts.file_analysis_prompt.format(
                            filename=filename,
                            file_type="media",
                            content="[Media content attached]",
                            query=state.request.query
                        )

                        agent = create_vlm_agent(
                            name="kb_vlm_prompt_agent",
                            vlm_mode=vlm_mode,
                            model_name="gemini-2.5-flash",
                            chat_api_url=state.request.chat_api_url,
                            parser_type="text",
                            additional_params={input_key: file_path}
                        )

                        temp_state = MainState(request=state.request)
                        temp_state.temp_data["kb_vlm_prompt"] = vlm_prompt

                        res_state = await agent.execute(temp_state)
                        analysis_result = _extract_text_result(res_state, "kb_vlm_prompt_agent")
                        if analysis_result:
                            raw_content = "[Media Content Processed by VLM]"
                        else:
                            analysis_result = "[VLM returned no content]"
                    except Exception as e:
                        log.error(f"VLM analysis failed for {filename}: {e}")
                        analysis_result = f"[VLM Analysis Error: {e}]"

                else:
                    # Use Simple Agent for Text
                    # Only if content is not empty or error
                    if raw_content and not raw_content.startswith("[Error"):
                        try:
                            # Prepare Prompt
                            # Limit raw_content size to avoid context overflow if huge (simple truncation)
                            truncated_content = raw_content[:50000]  # 50k char limit rough guard

                            analysis_prompt = QaAgentPrompts.file_analysis_prompt.format(
                                filename=filename,
                                file_type=file_type,
                                content=truncated_content,
                                query=state.request.query
                            )

                            agent = create_agent(
                                name="kb_prompt_agent",
                                model_name=state.request.model,
                                chat_api_url=state.request.chat_api_url,
                                temperature=0.3,
                                parser_type="text"
                            )

                            temp_state = MainState(request=state.request)
                            res_state = await agent.execute(temp_state, prompt=analysis_prompt)

                            analysis_result = _extract_text_result(res_state, "kb_prompt_agent")
                            if not analysis_result:
                                analysis_result = "[LLM Analysis Failed]"
                        except Exception as e:
                            log.error(f"Text analysis failed for {filename}: {e}")
                            analysis_result = f"[Text Analysis Error: {e}]"
                    else:
                        analysis_result = raw_content  # Pass through error or empty

            except Exception as e:
                 analysis_result = f"[Analysis Error: {e}]"
            
            return {
                "filename": filename,
                "analysis": analysis_result,
                "content": raw_content[:1000] + "..." if len(raw_content) > 1000 else raw_content # Store brief raw content for debug
            }

        # Run in parallel
        tasks = [process_file(f) for f in files_to_process]
        results = await asyncio.gather(*tasks)
        
        state.file_analyses = results
        return state

    def _try_rag_retrieve(state: IntelligentQAState) -> None:
        """若配置了 vector_store_base_dir 且索引存在，按 query 检索 Top-K 片段并写入 state.retrieved_chunks。"""
        base_dir = getattr(state.request, "vector_store_base_dir", None) or ""
        log.info(f"[_try_rag_retrieve] base_dir={base_dir}")
        log.info(f"[_try_rag_retrieve] file_ids={state.request.file_ids}")
        log.info(f"[_try_rag_retrieve] query={state.request.query[:100] if state.request.query else None}...")

        if not base_dir:
            log.warning("[_try_rag_retrieve] Skipped: no base_dir")
            return
        if not state.request.file_ids:
            log.warning("[_try_rag_retrieve] Skipped: no file_ids")
            return
        if not state.request.query:
            log.warning("[_try_rag_retrieve] Skipped: no query")
            return

        base_path = Path(base_dir)
        if not base_path.exists():
            log.warning(f"[_try_rag_retrieve] Skipped: base_dir does not exist: {base_dir}")
            return

        try:
            from workflow_engine.toolkits.ragtool.vector_store_tool import VectorStoreManager
            manager = VectorStoreManager(base_dir=base_dir)
            if manager.index is None or manager.index.ntotal == 0:
                log.warning(f"[_try_rag_retrieve] Skipped: index is empty (ntotal={manager.index.ntotal if manager.index else 0})")
                return

            # 从 manifest 中按「选中文件路径」解析出 kb file_ids
            manifest_files = manager.manifest.get("files", []) or []
            log.info(f"[_try_rag_retrieve] Manifest has {len(manifest_files)} files")

            local_paths = {Path(p).resolve() for p in state.request.file_ids}
            log.info(f"[_try_rag_retrieve] Looking for local_paths: {local_paths}")

            file_ids = []
            for f in manifest_files:
                orig = f.get("original_path") or ""
                if not orig:
                    continue
                try:
                    resolved = Path(orig).resolve()
                    if resolved in local_paths:
                        file_ids.append(f.get("id"))
                        log.info(f"[_try_rag_retrieve] ✓ Matched: {orig} (id={f.get('id')})")
                except Exception as e:
                    log.debug(f"[_try_rag_retrieve] Failed to resolve {orig}: {e}")

            if not file_ids:
                log.warning("[_try_rag_retrieve] No file_ids matched in manifest, will search all files")
                file_ids = None

            log.info(f"[_try_rag_retrieve] Searching with file_ids={file_ids}, top_k={RAG_TOP_K}")

            results = manager.search(
                query=state.request.query,
                top_k=RAG_TOP_K,
                file_ids=file_ids,
            )
            state.retrieved_chunks = results
            log.info(f"RAG 检索到 {len(results)} 个片段")
        except Exception as e:
            log.warning(f"RAG 检索跳过: {e}")
            state.retrieved_chunks = []

    def _format_history(history: List[Dict[str, str]]) -> str:
        """将 state.request.history 格式化为对话历史字符串，最多保留最近 MAX_HISTORY_TURNS 轮。"""
        if not history:
            return ""
        recent = history[-MAX_HISTORY_TURNS * 2:]
        lines = []
        for msg in recent:
            role = msg.get("role", "user")
            content = msg.get("content", "")
            if role == "user":
                lines.append(f"User: {content}")
            else:
                lines.append(f"Assistant: {content}")
        return "\n".join(lines)

    def _build_doc_context(state: IntelligentQAState) -> str:
        """构建文档上下文字符串：RAG 片段优先填充，file_analyses 补充，共享 MAX_DOC_CONTEXT_CHARS 预算。
        编号按 state.request.file_ids 顺序（与前端左侧来源面板一致），映射表存入 state.source_mapping。"""
        parts: List[str] = []
        total = 0

        # 按 request.files 顺序建立 filename -> 编号 映射
        file_index_map: Dict[str, int] = {}
        source_mapping: Dict[int, str] = {}
        for idx, fpath in enumerate(state.request.file_ids, 1):
            name = Path(fpath).name
            file_index_map[name] = idx
            source_mapping[idx] = name
        state.source_mapping = source_mapping

        # 1) RAG chunks（精准片段，优先填充）
        if state.retrieved_chunks:
            manifest_files: Dict[str, str] = {}
            try:
                from workflow_engine.toolkits.ragtool.vector_store_tool import VectorStoreManager
                base_dir = getattr(state.request, "vector_store_base_dir", None) or ""
                if base_dir:
                    mgr = VectorStoreManager(base_dir=base_dir)
                    for f in (mgr.manifest.get("files") or []):
                        if f.get("id"):
                            manifest_files[f["id"]] = Path(f.get("original_path", "")).name or f["id"]
            except Exception:
                pass
            for item in state.retrieved_chunks:
                content = (item.get("content") or "").strip()
                if not content:
                    continue
                fid = item.get("source_file_id", "")
                name = manifest_files.get(fid, fid)
                block = f"--- [{name}] ---\n{content}\n\n"
                if total + len(block) > MAX_DOC_CONTEXT_CHARS:
                    break
                parts.append(block)
                total += len(block)

        # 2) file_analyses（非入库文件的全文分析，补充剩余预算）
        if state.file_analyses:
            for item in state.file_analyses:
                fname = item.get('filename', '')
                block = f"--- Analysis of {fname} ---\n{item.get('analysis', '')}\n\n"
                if total + len(block) > MAX_DOC_CONTEXT_CHARS:
                    break
                parts.append(block)
                total += len(block)

        # 3) 附加编号来源映射表供 LLM 引用
        if source_mapping:
            mapping_lines = ["Sources:"]
            for idx in sorted(source_mapping.keys()):
                mapping_lines.append(f"[{idx}] {source_mapping[idx]}")
            parts.append("\n".join(mapping_lines) + "\n")

        return "".join(parts)

    async def chat_node(state: IntelligentQAState) -> IntelligentQAState:
        """
        Final synthesis: RAG 检索 + 文档上下文长度截断 + 对话历史。
        """
        _try_rag_retrieve(state)
        doc_context = _build_doc_context(state)
        history_str = _format_history(state.request.history)

        # 如果没有文件内容（file_analyses 和 retrieved_chunks 都为空），使用简化 prompt
        has_file_content = bool(state.file_analyses or state.retrieved_chunks)

        if not has_file_content:
            # 笔记 AI 辅助模式：query 本身包含完整上下文，不需要文件分析
            simple_prompt = state.request.query
            if history_str:
                simple_prompt = f"Conversation History:\n{history_str}\n\n{simple_prompt}"

            agent = create_agent(
                name="kb_prompt_agent",
                model_name=state.request.model,
                chat_api_url=state.request.chat_api_url,
                temperature=0.7,
                parser_type="text",
            )

            new_state = await agent.execute(state, prompt=simple_prompt)
            answer_text = _extract_text_result(new_state, "kb_prompt_agent")
            state.answer = answer_text or "Sorry, I couldn't generate an answer."
            return state

        # 标准模式：使用文件分析
        final_prompt = QaAgentPrompts.final_qa_prompt.format(
            query=state.request.query,
            file_analyses=doc_context,
            history=history_str,
        )

        agent = create_agent(
            name="kb_prompt_agent",
            model_name=state.request.model,
            chat_api_url=state.request.chat_api_url,
            temperature=0.7,
            parser_type="text",
        )

        new_state = await agent.execute(state, prompt=final_prompt)
        answer_text = _extract_text_result(new_state, "kb_prompt_agent")
        state.answer = answer_text or "Sorry, I couldn't generate an answer."

        return state

    nodes = {
        "_start_": _start_,
        "parallel_parse": parallel_parse_node,
        "chat": chat_node,
        "_end_": lambda s: s
    }

    edges = [
        ("_start_", "parallel_parse"),
        ("parallel_parse", "chat"),
        ("chat", "_end_")
    ]

    builder.add_nodes(nodes).add_edges(edges)
    return builder
