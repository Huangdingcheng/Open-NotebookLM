from __future__ import annotations
import os
import asyncio
from pathlib import Path
from typing import List, Dict, Any

import fitz  # PyMuPDF
from dataflow_agent.workflow.registry import register
from dataflow_agent.graphbuilder.graph_builder import GenericGraphBuilder
from dataflow_agent.logger import get_logger
from dataflow_agent.state import KBPodcastState, MainState
from dataflow_agent.agentroles import create_agent
from dataflow_agent.utils import get_project_root
import re
import wave
from dataflow_agent.toolkits.multimodaltool.req_tts import (
    generate_speech_bytes_async,
    split_tts_text
)

log = get_logger(__name__)

# Try importing office libraries
try:
    from docx import Document
except ImportError:
    Document = None

try:
    from pptx import Presentation
except ImportError:
    Presentation = None

@register("kb_podcast")
def create_kb_podcast_graph() -> GenericGraphBuilder:
    """
    Workflow for Knowledge Base Podcast Generation
    Steps:
    1. Parse uploaded files (PDF/Office)
    2. Generate podcast script using LLM
    3. Generate audio using TTS
    """
    builder = GenericGraphBuilder(state_model=KBPodcastState, entry_point="_start_")

    def _extract_text_result(state: MainState, role_name: str) -> str:
        try:
            result = state.agent_results.get(role_name, {}).get("results", {})
            if isinstance(result, dict):
                return result.get("text") or result.get("raw") or ""
            if isinstance(result, str):
                return result
        except Exception:
            return ""
        return ""

    def _start_(state: KBPodcastState) -> KBPodcastState:
        # Ensure request fields
        if not state.request.files:
            state.request.files = []

        # Initialize output directory
        if not state.result_path:
            project_root = get_project_root()
            import time
            ts = int(time.time())
            email = getattr(state.request, 'email', 'default')
            # Sanitize email for filesystem safety
            safe_email = re.sub(r'[^\w\-.]', '_', (email or 'default').replace('@', '_at_'))
            output_dir = project_root / "outputs" / "kb_outputs" / safe_email / f"{ts}_podcast"
            output_dir.mkdir(parents=True, exist_ok=True)
            state.result_path = str(output_dir)
        else:
            Path(state.result_path).mkdir(parents=True, exist_ok=True)

        state.file_contents = []
        state.podcast_script = ""
        state.audio_path = ""
        return state

    async def parse_files_node(state: KBPodcastState) -> KBPodcastState:
        """
        Parse all files and extract content
        """
        files = state.request.files
        if not files:
            state.file_contents = []
            return state

        async def process_file(file_path: str) -> Dict[str, Any]:
            file_path_obj = Path(file_path)
            filename = file_path_obj.name

            if not file_path_obj.exists():
                return {
                    "filename": filename,
                    "content": f"[Error: File not found {file_path}]"
                }

            suffix = file_path_obj.suffix.lower()
            raw_content = ""

            try:
                # PDF
                if suffix == ".pdf":
                    try:
                        doc = fitz.open(file_path)
                        text = ""
                        for page in doc:
                            text += page.get_text() + "\n"
                        raw_content = text
                    except Exception as e:
                        raw_content = f"[Error parsing PDF: {e}]"

                # Word
                elif suffix in [".docx", ".doc"]:
                    if Document is None:
                         raw_content = "[Error: python-docx not installed]"
                    else:
                        try:
                            doc = Document(file_path)
                            raw_content = "\n".join([p.text for p in doc.paragraphs])
                        except Exception as e:
                             raw_content = f"[Error parsing Docx: {e}]"

                # PPT
                elif suffix in [".pptx", ".ppt"]:
                    if Presentation is None:
                        raw_content = "[Error: python-pptx not installed]"
                    else:
                        try:
                            prs = Presentation(file_path)
                            text = ""
                            for i, slide in enumerate(prs.slides):
                                text += f"--- Slide {i+1} ---\n"
                                for shape in slide.shapes:
                                    if hasattr(shape, "text"):
                                        text += shape.text + "\n"
                            raw_content = text
                        except Exception as e:
                            raw_content = f"[Error parsing PPT: {e}]"

                else:
                    try:
                        with open(file_path, "r", encoding="utf-8") as f:
                            raw_content = f.read()
                    except:
                        raw_content = "[Unsupported file type]"

            except Exception as e:
                 raw_content = f"[Parse Error: {e}]"

            # Truncate content
            truncated_content = raw_content[:50000] if len(raw_content) > 50000 else raw_content

            return {
                "filename": filename,
                "content": truncated_content
            }

        # Run in parallel
        tasks = [process_file(f) for f in files]
        results = await asyncio.gather(*tasks)

        state.file_contents = results
        return state

    async def generate_script_node(state: KBPodcastState) -> KBPodcastState:
        """
        Generate podcast script using LLM
        """
        if not state.file_contents:
            state.podcast_script = "No content available for podcast generation."
            return state

        # Format file contents
        contents_str = ""
        for item in state.file_contents:
            contents_str += f"=== {item['filename']} ===\n{item['content']}\n\n"

        # Podcast script prompt
        language = state.request.language
        mode = getattr(state.request, "podcast_mode", "monologue")

        # Qwen TTS 只支持单人叙述，强制使用 monologue 模式
        use_local_tts = os.getenv("USE_LOCAL_TTS", "0").strip().lower() in ("1", "true", "yes")
        tts_engine = os.getenv("TTS_ENGINE", "qwen").strip().lower()
        if use_local_tts and tts_engine == "qwen" and mode == "dialog":
            log.info("[Podcast] Qwen TTS 不支持对话模式，转换为单人叙述")
            mode = "monologue"

        if mode == "dialog":
            if language == "zh":
                prompt = f"""你是一位专业的知识播客制作人。基于以下资料，生成一段5-10分钟的双人对话播客脚本。

角色设定：
- 说话人1（主持人）：引导话题、提问、总结要点，语气亲切自然
- 说话人2（嘉宾）：深入讲解、举例说明、回答问题，语气专业但不失幽默

对话要求：
1. 开场：主持人简短介绍主题，嘉宾用一个有趣的例子或观点吸引听众
2. 展开：通过问答形式逐步深入，主持人适时提出听众可能关心的问题
3. 互动：两人自然交流，可以有"确实是这样"、"这个很有意思"等回应
4. 类比：用生活化的比喻解释复杂概念
5. 节奏：避免一人长篇大论，保持一问一答或短段交流
6. 结尾：主持人总结核心要点，嘉宾给出实用建议或展望

格式要求（严格遵守）：
- 每行一个角色，格式为"[S1]..."（主持人）或"[S2]..."（嘉宾）
- 每段对话控制在1-3句话，保持自然节奏
- 避免书面语，使用口语化表达

资料内容：
{contents_str}

请生成播客脚本："""
            else:
                prompt = f"""You are a professional podcast producer. Based on the following materials, create a 5-10 minute two-person dialogue podcast script.

Character Setup:
- Speaker 1 (Host): Guides topics, asks questions, summarizes key points with a friendly tone
- Speaker 2 (Guest): Provides in-depth explanations, examples, answers questions with expertise and humor

Dialogue Requirements:
1. Opening: Host briefly introduces the topic, Guest hooks listeners with an interesting example
2. Development: Gradually dive deeper through Q&A, Host asks questions listeners might have
3. Interaction: Natural conversation flow with responses like "Exactly", "That's fascinating"
4. Analogies: Use relatable metaphors to explain complex concepts
5. Pacing: Avoid long monologues, maintain back-and-forth rhythm
6. Closing: Host summarizes key points, Guest offers practical advice or insights

Format Requirements (strictly follow):
- One speaker per line: "[S1]..." (Host) or "[S2]..." (Guest)
- Keep each dialogue turn to 1-3 sentences for natural flow
- Use conversational language, avoid formal writing style

Materials:
{contents_str}

Generate the podcast script:"""
        else:
            if language == "zh":
                prompt = f"""你是一位专业的知识播客主播。基于以下资料，生成一段5-10分钟的知识播客口播稿。

重要格式要求（必须严格遵守）：
1. 输出纯文本，不要使用任何markdown格式（不要**、#、---、*、列表符号等）
2. 不要添加舞台指示（不要【】、()等标注音乐、音效的内容）
3. 不要添加角色标签（不要"主播："、"旁白："等）
4. 直接输出可以朗读的连贯文本，就像你在对着麦克风说话

内容要求：
1. 口语化、生动有趣，像和朋友聊天一样自然
2. 结构清晰：开场问候 → 引入主题 → 核心内容讲解 → 总结收尾
3. 使用生活化的比喻和例子帮助理解
4. 适当加入互动性语言（"你可能会想..."、"想象一下..."）
5. 避免书面语和专业术语堆砌，用简单的话解释复杂概念
6. 段落之间用自然的过渡语连接

资料内容：
{contents_str}

请直接生成可以朗读的播客口播稿（纯文本，无格式）："""
            else:
                prompt = f"""You are a professional podcast host. Based on the following materials, create a 5-10 minute podcast narration script.

Critical Format Requirements (must strictly follow):
1. Output plain text only, no markdown formatting (no **, #, ---, *, bullet points, etc.)
2. No stage directions (no [], () for music, sound effects annotations)
3. No speaker labels (no "Host:", "Narrator:", etc.)
4. Output continuous readable text as if you're speaking into a microphone

Content Requirements:
1. Conversational and engaging, like chatting with a friend
2. Clear structure: greeting → topic introduction → core content → closing summary
3. Use relatable analogies and examples
4. Include interactive language ("You might wonder...", "Imagine...")
5. Avoid jargon, explain complex concepts in simple terms
6. Use natural transitions between sections

Materials:
{contents_str}

Generate the podcast narration script (plain text, no formatting):"""

        try:
            agent = create_agent(
                name="kb_prompt_agent",
                model_name=state.request.model,
                chat_api_url=state.request.chat_api_url,
                temperature=0.7,
                parser_type="text"
            )

            temp_state = MainState(request=state.request)
            res_state = await agent.execute(temp_state, prompt=prompt)

            state.podcast_script = _extract_text_result(res_state, "kb_prompt_agent") or "[Script generation failed]"
        except Exception as e:
            log.error(f"Script generation failed: {e}")
            state.podcast_script = f"[Script generation error: {e}]"

        # Save script to file
        try:
            script_path = Path(state.result_path) / "script.txt"
            script_path.write_text(state.podcast_script, encoding="utf-8")
        except Exception as e:
            log.error(f"Failed to save script: {e}")

        return state

    async def generate_audio_node(state: KBPodcastState) -> KBPodcastState:
        """
        Generate audio using TTS
        """
        if not state.podcast_script or (state.podcast_script.startswith("[") and not state.podcast_script.startswith("[S1]") and not state.podcast_script.startswith("[S2]")):
            state.audio_path = ""
            return state

        try:
            audio_path = str(Path(state.result_path) / "podcast.wav")
            mode = getattr(state.request, "podcast_mode", "monologue")

            # Debug logging
            use_local_tts = os.getenv("USE_LOCAL_TTS", "0").strip().lower() in ("1", "true", "yes")
            tts_engine = os.getenv("TTS_ENGINE", "qwen").strip().lower()
            log.info(f"[TTS DEBUG] mode={mode}, use_local_tts={use_local_tts}, tts_engine={tts_engine}")

            # Qwen TTS 只支持单人叙述，强制使用 monologue 模式
            if use_local_tts and tts_engine == "qwen":
                if mode == "dialog":
                    log.info("[TTS] Qwen TTS 不支持对话模式，转换为单人叙述")
                    mode = "monologue"

            # 本地 FireRedTTS2：分块生成对话
            if use_local_tts and mode == "dialog" and tts_engine == "firered":
                log.info("[TTS] 使用本地 FireRedTTS2 分块生成对话")

                # 检查脚本是否已经是 [S1]/[S2] 格式
                if "[S1]" in state.podcast_script or "[S2]" in state.podcast_script:
                    converted_text = state.podcast_script
                else:
                    language = state.request.language
                    speaker_a = "主持人" if language == "zh" else "Host"
                    speaker_b = "嘉宾" if language == "zh" else "Guest"
                    converted_script = []
                    for line in state.podcast_script.splitlines():
                        line = line.strip()
                        if not line:
                            continue
                        if line.startswith(speaker_a + ":") or line.startswith(speaker_a + "："):
                            converted_script.append("[S1]" + line.split(":", 1)[-1].split("：", 1)[-1].strip())
                        elif line.startswith(speaker_b + ":") or line.startswith(speaker_b + "："):
                            converted_script.append("[S2]" + line.split(":", 1)[-1].split("：", 1)[-1].strip())
                        else:
                            converted_script.append(line)
                    converted_text = "\n".join(converted_script)

                # 分块处理（每块最多6行对话，避免超过模型限制）
                lines = [l for l in converted_text.split("\n") if l.strip()]
                chunks = []
                for i in range(0, len(lines), 6):
                    chunks.append("\n".join(lines[i:i+6]))

                log.info(f"[TTS] 分为 {len(chunks)} 块生成")
                audio_chunks = []
                for i, chunk in enumerate(chunks):
                    log.info(f"[TTS] 生成第 {i+1}/{len(chunks)} 块")
                    audio_bytes = await generate_speech_bytes_async(
                        text=chunk,
                        api_url=state.request.chat_api_url,
                        api_key=state.request.api_key,
                        model=state.request.tts_model,
                        voice_name="Speaker1",
                    )
                    audio_chunks.append(audio_bytes)

                with open(audio_path, "wb") as f:
                    f.write(b"".join(audio_chunks))
                state.audio_path = audio_path
                log.info(f"[TTS] 音频已保存: {audio_path}")
                return state
            max_chars = 1500
            concurrency = 4

            segments = []
            if mode == "dialog":
                language = state.request.language
                speaker_a = "主持人" if language == "zh" else "Host"
                speaker_b = "嘉宾" if language == "zh" else "Guest"
                speaker_map = {
                    speaker_a.lower(): "A",
                    speaker_b.lower(): "B",
                    "a": "A",
                    "b": "B",
                    "speaker a": "A",
                    "speaker b": "B",
                    "角色a": "A",
                    "角色b": "B",
                    "主播": "A",
                    "嘉宾": "B",
                }
                pattern = re.compile(r"^\s*([^:：]{1,20})\s*[:：]\s*(.+)$")
                current_speaker = "A"
                for raw_line in state.podcast_script.splitlines():
                    line = raw_line.strip()
                    if not line:
                        continue
                    m = pattern.match(line)
                    if m:
                        label = m.group(1).strip().lower()
                        content = m.group(2).strip()
                        mapped = speaker_map.get(label)
                        if mapped:
                            current_speaker = mapped
                        if content:
                            segments.append({"speaker": current_speaker, "text": content})
                        continue
                    # No label, append to current speaker
                    if segments and segments[-1]["speaker"] == current_speaker:
                        segments[-1]["text"] = f"{segments[-1]['text']} {line}"
                    else:
                        segments.append({"speaker": current_speaker, "text": line})

                expanded = []
                for seg in segments:
                    for chunk in split_tts_text(seg["text"], max_chars):
                        expanded.append({
                            "speaker": seg["speaker"],
                            "text": chunk
                        })
                segments = expanded
            else:
                for chunk in split_tts_text(state.podcast_script, max_chars):
                    segments.append({"speaker": "A", "text": chunk})

            if not segments:
                raise RuntimeError("No valid TTS segments generated from script")

            sem = asyncio.Semaphore(concurrency)

            async def _run(seg):
                voice = state.request.voice_name if seg["speaker"] == "A" else state.request.voice_name_b
                async with sem:
                    return await generate_speech_bytes_async(
                        text=seg["text"],
                        api_url=state.request.chat_api_url,
                        api_key=state.request.api_key,
                        model=state.request.tts_model,
                        voice_name=voice,
                    )

            async def _run_no_sem(seg):
                voice = state.request.voice_name if seg["speaker"] == "A" else state.request.voice_name_b
                return await generate_speech_bytes_async(
                    text=seg["text"],
                    api_url=state.request.chat_api_url,
                    api_key=state.request.api_key,
                    model=state.request.tts_model,
                    voice_name=voice,
                )

            async def _run_with_retry(seg, attempts=3, base_delay=0.8, use_sem=True):
                last_err = None
                for i in range(attempts):
                    try:
                        if use_sem:
                            return await _run(seg)
                        return await _run_no_sem(seg)
                    except Exception as e:
                        last_err = e
                        if i < attempts - 1:
                            await asyncio.sleep(base_delay * (i + 1))
                        continue
                raise last_err

            tasks = [asyncio.create_task(_run_with_retry(seg, attempts=2, use_sem=True)) for seg in segments]
            results = await asyncio.gather(*tasks, return_exceptions=True)

            failed_indices = [i for i, r in enumerate(results) if isinstance(r, Exception)]
            if failed_indices:
                log.warning(f"TTS retry sequentially for {len(failed_indices)} failed segment(s)")
                for i in failed_indices:
                    results[i] = await _run_with_retry(segments[i], attempts=3, use_sem=False)

            audio_chunks = results

            os.makedirs(os.path.dirname(os.path.abspath(audio_path)), exist_ok=True)
            with wave.open(audio_path, "wb") as wav_file:
                wav_file.setnchannels(1)        # 1 Channel
                wav_file.setsampwidth(2)        # 16 bit = 2 bytes
                wav_file.setframerate(24000)    # 24kHz
                wav_file.writeframes(b"".join(audio_chunks))

            state.audio_path = audio_path
            log.info(f"Audio generated successfully: {audio_path}")
        except Exception as e:
            import traceback
            log.error(f"Audio generation failed: {e}")
            log.error(f"Traceback: {traceback.format_exc()}")
            err_str = str(e)
            if "503" in err_str or "model_not_found" in err_str or "model not found" in err_str.lower():
                state.audio_path = (
                    "[TTS 模型不可用：当前 API 不支持所选 TTS 模型（如 gemini-2.5-pro-preview-tts）。"
                    "请到「播客」设置中更换 TTS 模型，或使用支持该模型的 API 服务商。]"
                )
            else:
                state.audio_path = f"[Audio generation error: {e}]"

        return state

    nodes = {
        "_start_": _start_,
        "parse_files": parse_files_node,
        "generate_script": generate_script_node,
        "generate_audio": generate_audio_node,
        "_end_": lambda s: s
    }

    edges = [
        ("_start_", "parse_files"),
        ("parse_files", "generate_script"),
        ("generate_script", "generate_audio"),
        ("generate_audio", "_end_")
    ]

    builder.add_nodes(nodes).add_edges(edges)
    return builder
