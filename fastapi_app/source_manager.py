"""
Source manager — unified import / read for notebook sources.

Handles:
- Copying originals into the new directory layout
- Running MinerU for PDFs
- Generating unified markdown for every source type
- Reading back markdown / MinerU data for feature cards
- Fallback to legacy kb_data / kb_mineru paths
"""
from __future__ import annotations

import asyncio
import re
import shutil
import time
from dataclasses import dataclass, field
from pathlib import Path
from typing import List, Optional, Tuple

from workflow_engine.logger import get_logger
from workflow_engine.utils import get_project_root

from fastapi_app.notebook_paths import NotebookPaths

log = get_logger(__name__)


@dataclass
class SourceInfo:
    stem: str
    original_path: Path
    markdown_path: Optional[Path] = None
    mineru_path: Optional[Path] = None
    file_type: str = ""  # pdf, md, docx, pptx, url, text


class SourceManager:
    """Manage sources for one notebook."""

    def __init__(self, paths: NotebookPaths):
        self.paths = paths

    # ------------------------------------------------------------------
    # Import
    # ------------------------------------------------------------------

    async def import_file(self, file_path: Path, filename: str) -> SourceInfo:
        """
        Import a local file into the notebook source tree.
        1. Copy to sources/{stem}/original/
        2. PDF → run MinerU → sources/{stem}/mineru/
        3. Generate unified markdown → sources/{stem}/markdown/{stem}.md
        """
        stem = Path(filename).stem
        ext = Path(filename).suffix.lower()

        # 1) Copy original
        orig_dir = self.paths.source_original_dir(filename)
        orig_dir.mkdir(parents=True, exist_ok=True)
        dest = orig_dir / filename
        if file_path.resolve() != dest.resolve():
            shutil.copy2(str(file_path), str(dest))

        info = SourceInfo(
            stem=stem,
            original_path=dest,
            file_type=ext.lstrip(".") or "unknown",
        )

        # 2) PDF → MinerU
        if ext == ".pdf":
            mineru_dir = self.paths.source_mineru_dir(filename)
            mineru_dir.mkdir(parents=True, exist_ok=True)
            try:
                await self._run_mineru(dest, mineru_dir)
                info.mineru_path = mineru_dir
            except Exception as e:
                log.warning("[SourceManager] MinerU failed for %s: %s", filename, e)

        # 3) Unified markdown
        md_dir = self.paths.source_markdown_dir(filename)
        md_dir.mkdir(parents=True, exist_ok=True)
        md_path = md_dir / f"{stem}.md"
        md_text = self._generate_markdown(dest, ext, info.mineru_path)
        if md_text:
            md_path.write_text(md_text, encoding="utf-8")
            info.markdown_path = md_path

        return info

    async def import_text(self, content: str, title: str) -> SourceInfo:
        """Import plain text as a .md source."""
        safe = re.sub(r'[^\w\u4e00-\u9fff\s\-.]', "", (title or "").strip())
        safe = (safe or "text")[:80].strip() or "text"
        filename = f"{safe}_{int(time.time())}.md"
        stem = Path(filename).stem

        orig_dir = self.paths.source_original_dir(filename)
        orig_dir.mkdir(parents=True, exist_ok=True)
        dest = orig_dir / filename
        dest.write_text((content or "").strip(), encoding="utf-8")

        md_dir = self.paths.source_markdown_dir(filename)
        md_dir.mkdir(parents=True, exist_ok=True)
        md_path = md_dir / f"{stem}.md"
        shutil.copy2(str(dest), str(md_path))

        return SourceInfo(
            stem=stem,
            original_path=dest,
            markdown_path=md_path,
            file_type="text",
        )

    async def import_url(self, url: str, fetched_text: str, title: str = "") -> SourceInfo:
        """Import a URL source (text already fetched by caller)."""
        if not title:
            from urllib.parse import urlparse
            parsed = urlparse(url)
            title = (parsed.netloc or "web") + "_" + (parsed.path.strip("/") or "page")[:30]
        safe = re.sub(r'[^\w\u4e00-\u9fff\s\-.]', "", title)
        safe = (safe or "url")[:80].strip() or "url"
        filename = f"{safe}_{int(time.time())}.md"
        stem = Path(filename).stem

        orig_dir = self.paths.source_original_dir(filename)
        orig_dir.mkdir(parents=True, exist_ok=True)
        dest = orig_dir / filename
        dest.write_text(fetched_text.strip(), encoding="utf-8")

        md_dir = self.paths.source_markdown_dir(filename)
        md_dir.mkdir(parents=True, exist_ok=True)
        md_path = md_dir / f"{stem}.md"
        shutil.copy2(str(dest), str(md_path))

        return SourceInfo(
            stem=stem,
            original_path=dest,
            markdown_path=md_path,
            file_type="url",
        )

    # ------------------------------------------------------------------
    # Read helpers
    # ------------------------------------------------------------------

    def get_markdown(self, source_stem: str) -> str:
        """Read unified markdown for a source. Falls back to legacy paths."""
        # New path
        md = self._find_in_sources(source_stem, "markdown", "*.md")
        if md:
            return md

        # Fallback: try reading original directly (for .md files)
        orig = self._find_in_sources(source_stem, "original", "*.md")
        if orig:
            return orig

        return ""

    def get_mineru_md(self, source_stem: str) -> str:
        """Read MinerU markdown for a PDF source."""
        md = self._find_in_sources(source_stem, "mineru", "*.md")
        if md:
            return md

        # Fallback: search inside mineru/{stem}/auto/*.md
        mineru_dir = self.paths.sources_dir / source_stem / "mineru"
        if mineru_dir.exists():
            for pattern in ["*/auto/*.md", "*/hybrid_auto/*.md", "*/*.md"]:
                for f in mineru_dir.glob(pattern):
                    try:
                        return f.read_text(encoding="utf-8")
                    except Exception:
                        continue
        return ""

    def get_mineru_root(self, source_stem: str) -> Optional[Path]:
        """Return the MinerU auto/ directory path (for images etc.)."""
        mineru_dir = self.paths.sources_dir / source_stem / "mineru"
        if not mineru_dir.exists():
            return None
        # Look for auto/ or hybrid_auto/
        for sub in ("auto", "hybrid_auto"):
            candidate = mineru_dir / sub
            if candidate.is_dir():
                return candidate
        # Fallback: {stem}/auto inside mineru dir
        for child in sorted(mineru_dir.iterdir()):
            if child.is_dir():
                for sub in ("auto", "hybrid_auto"):
                    c2 = child / sub
                    if c2.is_dir():
                        return c2
                # Any dir with .md files
                if list(child.glob("*.md")):
                    return child
        return None

    def get_sam3_cache_dir(self, source_stem: str) -> Optional[Path]:
        """Return the SAM3 cache directory if it exists."""
        sam3_dir = self.paths.sources_dir / source_stem / "sam3"
        return sam3_dir if sam3_dir.exists() else None

    def get_sam3_results(self, source_stem: str) -> Optional[dict]:
        """Read cached SAM3 results JSON for a source."""
        sam3_dir = self.get_sam3_cache_dir(source_stem)
        if not sam3_dir:
            return None
        json_path = sam3_dir / "sam3_results.json"
        if not json_path.exists():
            return None
        try:
            import json
            return json.loads(json_path.read_text(encoding="utf-8"))
        except Exception:
            return None

    def get_sam3_elements(self, source_stem: str) -> Optional[list]:
        """Read cached SAM3 drawio elements for a source."""
        sam3_dir = self.get_sam3_cache_dir(source_stem)
        if not sam3_dir:
            return None
        json_path = sam3_dir / "drawio_elements.json"
        if not json_path.exists():
            return None
        try:
            import json
            return json.loads(json_path.read_text(encoding="utf-8"))
        except Exception:
            return None

    def ensure_sam3_dir(self, source_stem: str) -> Path:
        """Create and return the SAM3 cache directory."""
        sam3_dir = self.paths.sources_dir / source_stem / "sam3"
        sam3_dir.mkdir(parents=True, exist_ok=True)
        return sam3_dir

    def get_all_markdowns(self) -> List[Tuple[str, str]]:
        """Return [(stem, markdown_text), ...] for all sources."""
        results: List[Tuple[str, str]] = []
        sources_dir = self.paths.sources_dir
        if not sources_dir.exists():
            return results
        for src_dir in sorted(sources_dir.iterdir()):
            if not src_dir.is_dir():
                continue
            stem = src_dir.name
            md = self.get_markdown(stem)
            if md.strip():
                results.append((stem, md))
        return results

    def get_original_path(self, source_stem: str) -> Optional[Path]:
        """Return the original file path for a source."""
        orig_dir = self.paths.sources_dir / source_stem / "original"
        if not orig_dir.exists():
            return None
        for f in orig_dir.iterdir():
            if f.is_file():
                return f
        return None

    def list_sources(self) -> List[SourceInfo]:
        """List all sources in this notebook."""
        results: List[SourceInfo] = []
        sources_dir = self.paths.sources_dir
        if not sources_dir.exists():
            return results
        for src_dir in sorted(sources_dir.iterdir()):
            if not src_dir.is_dir():
                continue
            stem = src_dir.name
            orig = self.get_original_path(stem)
            if not orig:
                continue
            ext = orig.suffix.lower().lstrip(".")
            md_dir = src_dir / "markdown"
            md_path = None
            if md_dir.exists():
                mds = list(md_dir.glob("*.md"))
                if mds:
                    md_path = mds[0]
            mineru_path = src_dir / "mineru" if (src_dir / "mineru").exists() else None
            results.append(SourceInfo(
                stem=stem,
                original_path=orig,
                markdown_path=md_path,
                mineru_path=mineru_path,
                file_type=ext,
            ))
        return results

    # ------------------------------------------------------------------
    # Internal
    # ------------------------------------------------------------------

    async def _run_mineru(self, pdf_path: Path, output_dir: Path) -> None:
        """Run MinerU on a PDF file."""
        from workflow_engine.toolkits.multimodaltool.mineru_tool import run_mineru_pdf_extract
        await asyncio.to_thread(
            run_mineru_pdf_extract,
            str(pdf_path),
            str(output_dir),
            "modelscope",
            None,
            "pipeline",
        )

    def _generate_markdown(
        self, file_path: Path, ext: str, mineru_dir: Optional[Path]
    ) -> str:
        """Generate unified markdown from a source file."""
        # PDF: copy MinerU's .md output
        if ext == ".pdf" and mineru_dir:
            for pattern in ["*.md", "*/auto/*.md", "*/hybrid_auto/*.md"]:
                for f in mineru_dir.glob(pattern):
                    try:
                        return f.read_text(encoding="utf-8")
                    except Exception:
                        continue
            # Fallback: PyMuPDF extraction
            return self._extract_text_pymupdf(file_path)

        if ext == ".pdf":
            return self._extract_text_pymupdf(file_path)

        # MD / TXT: direct copy
        if ext in (".md", ".markdown", ".txt"):
            try:
                return file_path.read_text(encoding="utf-8", errors="replace")
            except Exception:
                return ""

        # DOCX
        if ext in (".docx", ".doc"):
            return self._extract_text_docx(file_path)

        # PPTX
        if ext in (".pptx", ".ppt"):
            return self._extract_text_pptx(file_path)

        # Fallback: try reading as text
        try:
            return file_path.read_text(encoding="utf-8", errors="replace")
        except Exception:
            return ""

    def _extract_text_pymupdf(self, path: Path) -> str:
        try:
            import fitz
            doc = fitz.open(str(path))
            text = "\n\n".join(page.get_text() for page in doc)
            doc.close()
            return text
        except Exception as e:
            log.warning("[SourceManager] PyMuPDF extraction failed: %s", e)
            return ""

    def _extract_text_docx(self, path: Path) -> str:
        try:
            from docx import Document
            doc = Document(str(path))
            return "\n".join(p.text for p in doc.paragraphs)
        except Exception as e:
            log.warning("[SourceManager] docx extraction failed: %s", e)
            return ""

    def _extract_text_pptx(self, path: Path) -> str:
        try:
            from pptx import Presentation
            prs = Presentation(str(path))
            parts = []
            for i, slide in enumerate(prs.slides):
                parts.append(f"--- Slide {i+1} ---")
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        parts.append(shape.text)
            return "\n".join(parts)
        except Exception as e:
            log.warning("[SourceManager] pptx extraction failed: %s", e)
            return ""

    def _find_in_sources(self, source_stem: str, subdir: str, pattern: str) -> str:
        """Find and read the first matching file in sources/{stem}/{subdir}/."""
        d = self.paths.sources_dir / source_stem / subdir
        if not d.exists():
            return ""
        for f in d.glob(pattern):
            try:
                return f.read_text(encoding="utf-8", errors="replace")
            except Exception:
                continue
        return ""
